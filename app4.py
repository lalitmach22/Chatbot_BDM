import os
import streamlit as st
import re
import json
from datetime import datetime, timedelta
import pytz
from docx import Document  # For .docx
import pandas as pd
from bs4 import BeautifulSoup
import json
import yaml
from pptx import Presentation
from zipfile import ZipFile
import mimetypes
from langchain_groq import ChatGroq
from langchain_community.document_loaders import PyPDFLoader
from langchain.chains import ConversationalRetrievalChain
from langchain.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
from supabase import create_client, Client

# Supabase credentials
url = "https://ycqqzosluyhqdwtnricr.supabase.co"
key = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InljcXF6b3NsdXlocWR3dG5yaWNyIiwicm9sZSI6ImFub24iLCJpYXQiOjE3MzQxODE5NjksImV4cCI6MjA0OTc1Nzk2OX0.Q-AmMeWuUkODkX27tg7YEN9bLvqz7v8qOHqpIucs_iw"
supabase: Client = create_client(url, key)

os.environ["GROQ_API_KEY"] = "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"

# Load the model
@st.cache_resource
def load_model():
    return ChatGroq(temperature=0.8, model="llama3-8b-8192")

def clean_text(text):
    # Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)
    
    # Fix broken sentences or words caused by line breaks
    text = re.sub(r'(?<=[a-zA-Z])\s(?=[a-zA-Z])', '', text)  # Remove single spaces in the middle of words

    # Standardize newlines for better formatting
    text = re.sub(r'\.\s+', '.\n', text)  # Add newlines after sentences
    text = re.sub(r'(?<=:)\s+', '\n', text)  # Add newlines after colons
    
    # Additional cleanup (if needed)
    text = text.strip()  # Remove leading and trailing whitespace

    return text

@st.cache_data
def load_hidden_documents(directory="hidden_docs"):
    """Load all supported file types from a directory and return their content."""
    all_texts = []

    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        mime_type, _ = mimetypes.guess_type(file_path)

        try:
            # Handle PDF files
            if filename.endswith(".pdf"):
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split()
                all_texts.extend([page.page_content for page in pages])

            # Handle Word files (.docx)
            elif filename.endswith(".docx"):
                doc = Document(file_path)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                all_texts.append(text)

            # Handle Text files (.txt)
            elif filename.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle Excel files (.xlsx and .xls)
            elif filename.endswith(('.xlsx', '.xls')):
                excel_data = pd.read_excel(file_path)
                text = excel_data.to_string(index=False)
                all_texts.append(text)

            # Handle CSV files (.csv)
            elif filename.endswith(".csv"):
                csv_data = pd.read_csv(file_path)
                text = csv_data.to_string(index=False)
                all_texts.append(text)

            # Handle Markdown files (.md)
            elif filename.endswith(".md"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle HTML files (.html, .htm)
            elif filename.endswith(('.html', '.htm')):
                with open(file_path, "r", encoding="utf-8") as file:
                    soup = BeautifulSoup(file, "html.parser")
                    all_texts.append(soup.get_text())

            # Handle JSON files (.json)
            elif filename.endswith(".json"):
                with open(file_path, "r", encoding="utf-8") as file:
                    data = json.load(file)
                    all_texts.append(json.dumps(data, indent=2))

            # Handle YAML files (.yaml, .yml)
            elif filename.endswith(('.yaml', '.yml')):
                with open(file_path, "r", encoding="utf-8") as file:
                    data = yaml.safe_load(file)
                    all_texts.append(json.dumps(data, indent=2))

            # Handle PowerPoint files (.pptx)
            elif filename.endswith(".pptx"):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_text.append(shape.text)
                    all_texts.append("\n".join(slide_text))

            # Handle ZIP files (.zip)
            elif filename.endswith(".zip"):
                with ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall("temp_extracted")
                    all_texts.extend(load_hidden_documents("temp_extracted"))

            # Handle Log files (.log)
            elif filename.endswith(".log"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle unknown file types (fallback to text-based reading)
            elif mime_type and mime_type.startswith("text"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

        except Exception as e:
            print(f"Failed to process {filename}: {e}")
    cleaned_texts = [clean_text(text) for text in all_texts]
    return cleaned_texts

@st.cache_data
def save_to_supabase(all_texts):
    """Save the list of documents to the Supabase 'all_texts' table."""
    for text in all_texts:
        data = {"all_texts": text}
        response = supabase.table("all_texts").insert(data).execute()
        
        # Check the response for success or failure
        if response.data:  # If the response contains data, the insert was successful
            print(f"Successfully saved: {text[:30]}...")
        else:  # If there is an error
            print(f"Failed to save text. Error")

# Load documents and save to Supabase
all_texts = load_hidden_documents()
save_to_supabase(all_texts)

# Create vector store
@st.cache_resource
def create_vector_store(document_texts):
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(document_texts, embedder)

# Get modification times for documents
@st.cache_data
def get_file_mod_times(directory):
    """Get the modification times of all files in the directory."""
    return {
        f: os.path.getmtime(os.path.join(directory, f))
        for f in os.listdir(directory)
        if os.path.isfile(os.path.join(directory, f))  # Ensure it's a file, not a directory
    }

# Reload vector store if needed
def reload_vector_store_if_needed():
    current_mod_times = get_file_mod_times("hidden_docs")

    # Check if file modifications have occurred
    if "file_mod_times" not in st.session_state or st.session_state["file_mod_times"] != current_mod_times:
        st.session_state["file_mod_times"] = current_mod_times
        document_texts = load_hidden_documents()
        vector_store = create_vector_store(document_texts)  # Create vector store
        st.session_state["vector_store"] = vector_store    # Save in session state
    else:
        # Retrieve from session state if already initialized
        vector_store = st.session_state.get("vector_store", None)

    # Return the vector store (even if None)
    return vector_store

# Load model and vector store
model = load_model()

# Initialize vector_store
vector_store = reload_vector_store_if_needed()

# If still None, raise an error to debug initialization
if vector_store is None:
    raise ValueError("Failed to initialize vector_store. Ensure hidden_docs folder and embeddings setup are correct.")

# Validate email
def is_valid_email(email):
    email_regex = r"^\d{2}f\d{7}@ds\.study\.iitm\.ac\.in$"
    return re.match(email_regex, email) is not None or email == "nitin@ee.iitm.ac.in"

# Save session to Supabase
def save_session(session_data):
    data = {"session_data": session_data}
    response = supabase.table("session_data").insert(data).execute()

# Set timeout for 30 minutes
session_start_time = datetime.now(pytz.timezone("Asia/Kolkata"))
session_timeout = timedelta(minutes=30)
if datetime.now(pytz.timezone("Asia/Kolkata")) - session_start_time > session_timeout:
    # Save session data
    save_session(session_data="Session expired.")
    st.session_state["session_expired"] = True

# Implement chat
def chat_with_bot(user_input):
    # Check for session expiry
    if "session_expired" in st.session_state and st.session_state["session_expired"]:
        st.warning("Session expired. Please start a new session.")
        return
    
    # Add user input to conversation history
    conversation_history = [{"role": "user", "content": user_input}]
    response = model.chat(messages=conversation_history)
    
    # Return the model's response
    return response

# Streamlit UI
def main():
    st.title("Session-Based Chatbot")
    st.text_area("Chat with the bot:", height=300)
    user_input = st.text_input("Enter your message:")

    if user_input:
        response = chat_with_bot(user_input)
        st.write(response['message'])
    
    if st.button("Download Session Data"):
        # Provide option to download session data as JSON file
        session_data = json.dumps(st.session_state)
        st.download_button("Download Session Data", data=session_data, file_name="session_data.json")

if __name__ == "__main__":
    main()